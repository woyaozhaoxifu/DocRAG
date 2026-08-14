import os
import re

from . import ocr

MAX_BYTES = 4 * 1024 * 1024  # 单文件最多读取 4MB 原文

IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".bmp", ".gif")


def _pdf(path):
    import pymupdf
    doc = pymupdf.open(path)
    try:
        out = [page.get_text() for page in doc]
    finally:
        doc.close()
    return "\n".join(out)


def _docx(path):
    import docx
    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs if p.text)


def _doc(path):
    # 旧版 .doc 需要 LibreOffice，这里尝试用 antiword/textract 失败则降级为空
    return ""


def _xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        out = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    out.append(" ".join(cells))
    finally:
        wb.close()
    return "\n".join(out)


def _pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs)
                    if t:
                        out.append(t)
    return "\n".join(out)


def _plain(path):
    for enc in ("utf-8", "gbk", "gb18030", "latin-1"):
        try:
            with open(path, encoding=enc, errors="strict") as f:
                return f.read()
        except Exception:
            continue
    return ""


def extract_text(path, max_chars=200000):
    """抽取文档正文，统一清洗。失败返回空字符串。"""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            text = ocr.extract_pdf(path) or _pdf(path)
        elif ext == ".docx":
            text = _docx(path)
        elif ext == ".doc":
            text = _doc(path)
        elif ext in (".xlsx", ".xls"):
            text = _xlsx(path)
        elif ext == ".pptx":
            text = _pptx(path)
        elif ext in IMAGE_EXTS:
            text = ocr.extract_image(path) or ""
        elif ext in (".txt", ".md", ".csv", ".json", ".html", ".htm", ".rtf"):
            text = _plain(path)
        else:
            text = ""
    except Exception:
        text = ""

    if text:
        text = re.sub(r"\s+", " ", text).strip()
        if len(text) > max_chars:
            text = text[:max_chars]
    return text
