"""端到端验证：解析 / 索引 / 关键词 / 搜索 / 归档 / 还原。"""
import os
import sys
import tempfile
import shutil

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from core.store import Store, initial_scan

PASS, FAIL = 0, 0


def check(name, cond):
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  [OK]   {name}")
    else:
        FAIL += 1
        print(f"  [FAIL] {name}")


def make_docs(d):
    files = {}
    # txt
    p = os.path.join(d, "人工智能报告.txt")
    with open(p, "w", encoding="utf-8") as f:
        f.write("人工智能与机器学习在医疗影像诊断中的应用日益广泛，深度学习模型提升了准确率。")
    files["txt"] = p
    # md
    p = os.path.join(d, "会议纪要.md")
    with open(p, "w", encoding="utf-8") as f:
        f.write("# 项目会议纪要\n讨论本地检索系统的架构设计与向量数据库选型。")
    files["md"] = p
    # docx
    import docx
    p = os.path.join(d, "方案.docx")
    doc = docx.Document()
    doc.add_paragraph("区块链技术在供应链金融中的落地方案与风控要点。")
    doc.save(p)
    files["docx"] = p
    # pdf
    import pymupdf
    p = os.path.join(d, "白皮书.pdf")
    pdf = pymupdf.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "新能源电池产业链分析：正极材料与固态电池技术路线。", fontname="china-s")
    pdf.save(p)
    pdf.close()
    files["pdf"] = p
    # xlsx
    import openpyxl
    p = os.path.join(d, "数据.xlsx")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["季度", "营收"])
    ws.append(["Q1", "1200万"])
    wb.save(p)
    files["xlsx"] = p
    # pptx
    from pptx import Presentation
    p = os.path.join(d, "演示.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "产品发布路线路演"
    prs.save(p)
    files["pptx"] = p
    return files


def main():
    tmp = tempfile.mkdtemp(prefix="docrag_test_")
    print("测试目录:", tmp)
    docs_dir = os.path.join(tmp, "docs")
    data_dir = os.path.join(tmp, "data")
    archive_root = os.path.join(tmp, "archive")
    os.makedirs(docs_dir, exist_ok=True)

    cfg = {
        "monitored_dirs": {"Test": docs_dir},
        "auto_archive": {"Test": False},
        "archive_root": archive_root,
        "data_dir": data_dir,
        "extensions": [".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
                       ".txt", ".md", ".csv", ".json", ".html", ".htm", ".rtf"],
        "max_text_chars": 200000,
        "settle_seconds": 1,
    }

    files = make_docs(docs_dir)

    store = Store(cfg)

    print("\n[1] 解析 + 索引")
    for k, p in files.items():
        ok = store.index_file(p)
        check(f"索引 {k} ({os.path.basename(p)})", ok)

    print("\n[2] 关键词抽取")
    res = store.search("人工智能", 10)
    check("搜索'人工智能'命中报告", any("人工智能报告" in (r["filename"] or "") for r in res))
    kw_ok = any(r.get("keywords") for r in res)
    check("返回结果含关键词", kw_ok)

    print("\n[3] 多格式内容检索")
    for term, expect in [("区块链", "方案"), ("固态电池", "白皮书"),
                         ("供应链金融", "方案"), ("营收", "数据"), ("路演", "演示")]:
        r = store.search(term, 10)
        hit = any(expect in (x["filename"] or "") for x in r)
        check(f"搜索'{term}'命中{expect}", hit)

    print("\n[4] 统计与近期")
    st = store.stats()
    check(f"已索引数量={st['total']}>=6", st["total"] >= 6)

    print("\n[5] 安全归档 + 还原")
    target = files["txt"]
    before = os.path.exists(target)
    arc = store.archive_file(target)
    check("归档成功", arc.get("ok"))
    check("源文件已移走", before and not os.path.exists(target))
    check("归档后存在于归档库", os.path.exists(arc["archive_path"]))
    # 归档后搜索仍能命中（路径变为归档路径）
    r2 = store.search("人工智能", 10)
    hit2 = any("人工智能报告" in (x["filename"] or "") for x in r2)
    check("归档后搜索仍命中", hit2)
    archived_flag = any(x.get("archived") for x in r2 if "人工智能报告" in (x["filename"] or ""))
    check("结果标记为已归档", archived_flag)
    # 还原
    res3 = store.restore(arc["archive_path"])
    check("还原成功", res3.get("ok"))
    check("还原回原路径", os.path.exists(target))
    check("归档库文件已移走", not os.path.exists(arc["archive_path"]))

    print("\n[6] initial_scan 只读扫描")
    cnt = initial_scan(store, cfg)
    check(f"initial_scan 索引数={cnt}>=6", cnt >= 6)

    store.conn.close()
    shutil.rmtree(tmp, ignore_errors=True)
    print(f"\n=== 结果: PASS={PASS} FAIL={FAIL} ===")
    sys.exit(1 if FAIL else 0)


if __name__ == "__main__":
    main()
